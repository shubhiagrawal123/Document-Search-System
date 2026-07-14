import json
import pandas as pd
import PyPDF2

from docx import Document
from pptx import Presentation


def read_pdf(path):

    text = ""

    with open(path, "rb") as file:

        reader = PyPDF2.PdfReader(file)

        for page in reader.pages:

            page_text = page.extract_text()

            if page_text:

                text += page_text + "\n"

    return text


def read_txt(path):

    with open(
        path,
        "r",
        encoding="utf-8",
        errors="ignore"
    ) as file:

        return file.read()


def read_docx(path):

    doc = Document(path)

    text = ""

    for para in doc.paragraphs:

        text += para.text + "\n"

    for table in doc.tables:

        for row in table.rows:

            for cell in row.cells:

                text += cell.text + "\n"

    return text


def read_csv(path):

    df = pd.read_csv(
        path,
        encoding="utf-8",
        on_bad_lines="skip"
    )

    return df.to_string()


def read_xlsx(path):

    text = ""

    excel_file = pd.ExcelFile(path)

    for sheet in excel_file.sheet_names:

        df = pd.read_excel(
            path,
            sheet_name=sheet
        )

        text += f"\n\n===== {sheet} =====\n"

        text += df.to_string()

    return text


def read_pptx(path):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text


def read_json(path):

    with open(
        path,
        "r",
        encoding="utf-8"
    ) as file:

        data = json.load(file)

    return json.dumps(
        data,
        indent=2
    )